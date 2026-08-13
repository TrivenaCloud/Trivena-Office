"""Brief → one-slide PPTX for TrivOffice (replaces Genspark tool_cli /slide_generate).

Produces several modern layouts (cover, cards, split image, two-column, big number)
from an LLM JSON plan (NVIDIA GLM / Gemini) + python-pptx. Not Genspark HTML quality, but far above a
single accent-bar + bullet list.
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
from typing import Any

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

_HEX = re.compile(r"#([0-9a-fA-F]{6})")
LAYOUTS = ("cover", "title_cards", "split_image", "two_column", "big_number", "bullets")


def _rgb(hex_color: str, default: str = "1A1A2E") -> RGBColor:
    m = _HEX.search(hex_color or "")
    h = m.group(1) if m else default
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _pick_colors(style_skill: str) -> dict[str, str]:
    bg, text, accent, card, muted = "F8FAFC", "0F172A", "2563EB", "FFFFFF", "64748B"
    skill = style_skill or ""
    for label, key in (
        ("Main background", "bg"),
        ("Main text", "text"),
        ("Primary accent", "accent"),
        ("Card background", "card"),
        ("Secondary accent", "muted"),
    ):
        m = re.search(rf"{label}[^#\n]*?(#[0-9a-fA-F]{{6}})", skill, re.I)
        if not m:
            continue
        val = m.group(1)
        if key == "bg":
            bg = val
        elif key == "text":
            text = val
        elif key == "accent":
            accent = val
        elif key == "card":
            card = val
        else:
            muted = val
    found = _HEX.findall(skill)
    if found and bg == "F8FAFC":
        if len(found) >= 1:
            bg = "#" + found[0]
        if len(found) >= 2:
            text = "#" + found[1]
        if len(found) >= 3:
            accent = "#" + found[2]
    return {
        "bg": bg if bg.startswith("#") else f"#{bg}",
        "text": text if text.startswith("#") else f"#{text}",
        "accent": accent if accent.startswith("#") else f"#{accent}",
        "card": card if card.startswith("#") else f"#{card}",
        "muted": muted if muted.startswith("#") else f"#{muted}",
    }


def _set_fill(shape: Any, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    shape.line.fill.background()


def _round_rect(slide: Any, left: Any, top: Any, width: Any, height: Any, fill: str, radius: float = 0.15) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _set_fill(shape, fill)
    # Softer corners
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    return shape


def _textbox(
    slide: Any,
    left: Any,
    top: Any,
    width: Any,
    height: Any,
    text: str,
    *,
    size: int,
    color: str,
    bold: bool = False,
    align: Any = PP_ALIGN.LEFT,
) -> Any:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.alignment = align
    return box


def _add_bullets(tf: Any, bullets: list[str], *, size: int, color: str) -> None:
    for i, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(size)
        para.font.color.rgb = _rgb(color)
        para.space_after = Pt(8)


def _llm_openai_endpoint() -> tuple[str, str, str]:
    """(api_key, base_url, model) for slide planning. Prefer NVIDIA NIM."""
    nvidia_key = os.environ.get("NVIDIA_API_KEY", "") or os.environ.get("NGC_API_KEY", "")
    if nvidia_key:
        base = os.environ.get("NVIDIA_OPENAI_BASE_URL", "https://integrate.api.nvidia.com/v1").rstrip(
            "/"
        )
        model = os.environ.get("NVIDIA_DEFAULT_MODEL", "z-ai/glm-5.2")
        return nvidia_key, base, model
    gemini_key = os.environ.get("GEMINI_API_KEY", "") or os.environ.get("GOOGLE_API_KEY", "")
    if gemini_key:
        base = os.environ.get(
            "GEMINI_OPENAI_BASE_URL",
            "https://generativelanguage.googleapis.com/v1beta/openai",
        ).rstrip("/")
        model = os.environ.get("GEMINI_DEFAULT_MODEL", "gemini-2.5-flash")
        return gemini_key, base, model
    return "", "", ""


async def _llm_structure(
    *,
    brief: str,
    title: str,
    style_skill: str,
    deck_context: dict[str, Any] | None,
) -> dict[str, Any]:
    api_key, base, model = _llm_openai_endpoint()
    if not api_key:
        lines = [ln.strip(" -•\t") for ln in brief.splitlines() if ln.strip()]
        return {
            "title": title or (lines[0][:80] if lines else "Slide"),
            "subtitle": "",
            "layout": "title_cards",
            "cards": [
                {"title": "Point", "body": b} for b in (lines[1:4] if len(lines) > 1 else lines[:3])
            ],
            "bullets": lines[1:6],
        }

    ctx = json.dumps(deck_context or {}, ensure_ascii=False)[:1500]
    prompt = (
        "You are a presentation designer. Reply with ONLY valid JSON (no markdown):\n"
        "{\n"
        '  "title": "short punchy title",\n'
        '  "subtitle": "optional one-liner",\n'
        '  "kicker": "optional eyebrow label",\n'
        '  "layout": "cover|title_cards|split_image|two_column|big_number|bullets",\n'
        '  "cards": [{"title":"...","body":"..."}],\n'
        '  "columns": [{"heading":"...","bullets":["..."]}],\n'
        '  "stat": {"value":"42%","label":"..."},\n'
        '  "bullets": ["..."],\n'
        '  "image_role": "hero|supporting|none"\n'
        "}\n"
        "Rules:\n"
        "- Pick the layout that best fits the brief (features→title_cards or two_column; "
        "cover/agenda→cover; KPI→big_number; narrative→split_image or bullets).\n"
        "- title_cards: 2–4 cards with short title+body.\n"
        "- two_column: exactly 2 columns with heading + 2–4 bullets each.\n"
        "- Cover: title + subtitle + optional kicker; few or no bullets.\n"
        "- Concrete facts from the brief only; no invented numbers.\n"
        "- Titles ≤ 8 words; card bodies ≤ 18 words.\n\n"
        f"TITLE HINT: {title}\n"
        f"STYLE: {style_skill[:1200]}\n"
        f"DECK CONTEXT: {ctx}\n"
        f"BRIEF:\n{brief[:6000]}"
    )
    url = f"{base}/chat/completions"
    payload = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "temperature": 0.5,
        "max_tokens": 1800,
    }
    async with httpx.AsyncClient(timeout=90.0) as client:
        resp = await client.post(
            url,
            headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
            json=payload,
        )
    if resp.status_code >= 400:
        raise RuntimeError(f"Slide plan LLM failed HTTP {resp.status_code}: {resp.text[:200]}")
    data = resp.json()
    content = (((data.get("choices") or [{}])[0].get("message") or {}).get("content") or "").strip()
    if content.startswith("```"):
        content = re.sub(r"^```(?:json)?\s*", "", content)
        content = re.sub(r"\s*```$", "", content)
    try:
        parsed = json.loads(content)
    except json.JSONDecodeError as e:
        raise RuntimeError(f"Slide plan LLM returned non-JSON: {content[:200]}") from e
    if not isinstance(parsed, dict):
        raise RuntimeError("Slide plan was not an object")
    return parsed


async def _fetch_image(url: str) -> bytes | None:
    try:
        async with httpx.AsyncClient(timeout=20.0, follow_redirects=True) as client:
            r = await client.get(url, headers={"User-Agent": "TrivOffice-SlideGen/1.0"})
        if r.status_code == 200 and (
            r.headers.get("content-type", "").startswith("image") or len(r.content) > 1000
        ):
            return r.content
    except Exception:
        return None
    return None


def _blank_slide(width_px: int, height_px: int) -> tuple[Presentation, Any, float, float]:
    w_in = max(8.0, min(20.0, (width_px or 1280) / 96.0))
    h_in = max(5.0, min(12.0, (height_px or 720) / 96.0))
    prs = Presentation()
    prs.slide_width = Inches(w_in)
    prs.slide_height = Inches(h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide, w_in, h_in


def _layout_cover(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str], img: bytes | None) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    # Accent panel left third
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w * 0.38), Inches(h))
    _set_fill(panel, c["accent"])
    kicker = str(s.get("kicker") or "").strip()
    title = str(s.get("title") or "Slide").strip()[:100]
    subtitle = str(s.get("subtitle") or "").strip()[:160]
    y = Inches(h * 0.28)
    if kicker:
        _textbox(slide, Inches(0.55), y, Inches(w * 0.3), Inches(0.4), kicker.upper(), size=12, color="FFFFFF", bold=True)
        y = Inches(h * 0.34)
    _textbox(slide, Inches(0.55), y, Inches(w * 0.3), Inches(2.2), title, size=36, color="FFFFFF", bold=True)
    if subtitle:
        _textbox(
            slide,
            Inches(0.55),
            Inches(h * 0.62),
            Inches(w * 0.3),
            Inches(1.2),
            subtitle,
            size=16,
            color="FFFFFF",
        )
    if img:
        try:
            slide.shapes.add_picture(io.BytesIO(img), Inches(w * 0.38), Emu(0), width=Inches(w * 0.62), height=Inches(h))
        except Exception:
            pass


def _layout_cards(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str]) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    # Top accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(0.12))
    _set_fill(stripe, c["accent"])
    title = str(s.get("title") or "Slide").strip()[:100]
    subtitle = str(s.get("subtitle") or "").strip()[:140]
    _textbox(slide, Inches(0.55), Inches(0.4), Inches(w - 1.1), Inches(0.7), title, size=30, color=c["text"], bold=True)
    if subtitle:
        _textbox(slide, Inches(0.55), Inches(1.05), Inches(w - 1.1), Inches(0.4), subtitle, size=14, color=c["muted"])

    cards = s.get("cards") if isinstance(s.get("cards"), list) else []
    cards = [x for x in cards if isinstance(x, dict)][:4]
    if not cards:
        bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
        cards = [{"title": f"{i+1}", "body": str(b)} for i, b in enumerate(bullets[:4])]
    n = max(1, len(cards))
    gap = 0.25
    margin = 0.55
    usable = w - 2 * margin - gap * (n - 1)
    cw = usable / n
    top_in = 1.55 if subtitle else 1.25
    ch = max(2.5, h - top_in - 0.55)
    for i, card in enumerate(cards):
        left = margin + i * (cw + gap)
        _round_rect(slide, Inches(left), Inches(top_in), Inches(cw), Inches(ch), c["card"], 0.12)
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top_in), Inches(cw), Inches(0.08)
        )
        _set_fill(accent_bar, c["accent"])
        ct = str(card.get("title") or "").strip()[:40]
        cb = str(card.get("body") or "").strip()[:120]
        _textbox(
            slide,
            Inches(left + 0.2),
            Inches(top_in + 0.35),
            Inches(cw - 0.4),
            Inches(0.6),
            ct,
            size=16,
            color=c["text"],
            bold=True,
        )
        _textbox(
            slide,
            Inches(left + 0.2),
            Inches(top_in + 1.05),
            Inches(cw - 0.4),
            Inches(ch - 1.4),
            cb,
            size=13,
            color=c["muted"],
        )


def _layout_split(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str], img: bytes | None) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    title = str(s.get("title") or "Slide").strip()[:100]
    _textbox(slide, Inches(0.55), Inches(0.4), Inches(w * 0.5), Inches(0.9), title, size=28, color=c["text"], bold=True)
    bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
    bullets = [str(b).strip() for b in bullets if str(b).strip()][:6]
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(w * 0.48), Inches(h - 2.0))
    _add_bullets(body.text_frame, bullets or [str(s.get("subtitle") or title)], size=16, color=c["text"])
    # Image panel with rounded card behind
    if img:
        _round_rect(slide, Inches(w * 0.55), Inches(0.55), Inches(w * 0.4), Inches(h - 1.1), c["card"], 0.1)
        try:
            slide.shapes.add_picture(
                io.BytesIO(img),
                Inches(w * 0.58),
                Inches(0.8),
                width=Inches(w * 0.34),
                height=Inches(h - 1.6),
            )
        except Exception:
            pass
    else:
        # Accent block instead of missing image
        _round_rect(slide, Inches(w * 0.55), Inches(0.55), Inches(w * 0.4), Inches(h - 1.1), c["accent"], 0.1)
        _textbox(
            slide,
            Inches(w * 0.6),
            Inches(h * 0.4),
            Inches(w * 0.3),
            Inches(1),
            str(s.get("kicker") or s.get("subtitle") or ""),
            size=18,
            color="FFFFFF",
            bold=True,
            align=PP_ALIGN.CENTER,
        )


def _layout_two_column(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str]) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(0.1))
    _set_fill(stripe, c["accent"])
    title = str(s.get("title") or "Slide").strip()[:100]
    _textbox(slide, Inches(0.55), Inches(0.35), Inches(w - 1.1), Inches(0.7), title, size=28, color=c["text"], bold=True)
    cols = s.get("columns") if isinstance(s.get("columns"), list) else []
    cols = [x for x in cols if isinstance(x, dict)][:2]
    if len(cols) < 2:
        bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
        mid = max(1, len(bullets) // 2)
        cols = [
            {"heading": "Kenmerken", "bullets": [str(b) for b in bullets[:mid]]},
            {"heading": "Voordelen", "bullets": [str(b) for b in bullets[mid:]]},
        ]
    for i, col in enumerate(cols):
        left = 0.55 + i * (w * 0.48)
        _round_rect(slide, Inches(left), Inches(1.25), Inches(w * 0.42), Inches(h - 1.9), c["card"], 0.1)
        _textbox(
            slide,
            Inches(left + 0.25),
            Inches(1.5),
            Inches(w * 0.36),
            Inches(0.5),
            str(col.get("heading") or f"Column {i+1}")[:40],
            size=18,
            color=c["accent"],
            bold=True,
        )
        bl = col.get("bullets") if isinstance(col.get("bullets"), list) else []
        bl = [str(b).strip() for b in bl if str(b).strip()][:5]
        box = slide.shapes.add_textbox(Inches(left + 0.25), Inches(2.15), Inches(w * 0.36), Inches(h - 3.0))
        _add_bullets(box.text_frame, bl or ["—"], size=14, color=c["text"])


def _layout_big_number(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str]) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    title = str(s.get("title") or "Slide").strip()[:80]
    stat = s.get("stat") if isinstance(s.get("stat"), dict) else {}
    value = str(stat.get("value") or "").strip() or "—"
    label = str(stat.get("label") or s.get("subtitle") or "").strip()[:80]
    _textbox(slide, Inches(0.7), Inches(0.5), Inches(w - 1.4), Inches(0.6), title, size=22, color=c["muted"], bold=True)
    _textbox(
        slide,
        Inches(0.7),
        Inches(h * 0.28),
        Inches(w - 1.4),
        Inches(1.8),
        value[:20],
        size=72,
        color=c["accent"],
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if label:
        _textbox(
            slide,
            Inches(1.5),
            Inches(h * 0.58),
            Inches(w - 3),
            Inches(1),
            label,
            size=20,
            color=c["text"],
            align=PP_ALIGN.CENTER,
        )
    bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
    bullets = [str(b).strip() for b in bullets if str(b).strip()][:3]
    if bullets:
        box = slide.shapes.add_textbox(Inches(1.5), Inches(h * 0.72), Inches(w - 3), Inches(1.4))
        _add_bullets(box.text_frame, bullets, size=14, color=c["muted"])


def _layout_bullets(slide: Any, w: float, h: float, s: dict[str, Any], c: dict[str, str]) -> None:
    _set_fill(
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(w), Inches(h)),
        c["bg"],
    )
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Inches(0.18), Inches(h))
    _set_fill(bar, c["accent"])
    title = str(s.get("title") or "Slide").strip()[:100]
    _textbox(slide, Inches(0.6), Inches(0.45), Inches(w - 1.2), Inches(0.9), title, size=30, color=c["text"], bold=True)
    bullets = s.get("bullets") if isinstance(s.get("bullets"), list) else []
    bullets = [str(b).strip() for b in bullets if str(b).strip()][:7]
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(w - 1.5), Inches(h - 2.1))
    _add_bullets(box.text_frame, bullets or [str(s.get("subtitle") or title)], size=18, color=c["text"])


def _build_pptx(
    structure: dict[str, Any],
    *,
    colors: dict[str, str],
    image_bytes: bytes | None,
    width_px: int,
    height_px: int,
) -> bytes:
    prs, slide, w, h = _blank_slide(width_px, height_px)
    layout = str(structure.get("layout") or "title_cards").strip().lower()
    if layout not in LAYOUTS:
        layout = "title_cards"
    # Prefer cards when Gemini forgot but gave cards
    if layout == "bullets" and isinstance(structure.get("cards"), list) and structure["cards"]:
        layout = "title_cards"
    if layout == "cover":
        _layout_cover(slide, w, h, structure, colors, image_bytes)
    elif layout == "title_cards":
        _layout_cards(slide, w, h, structure, colors)
    elif layout == "split_image":
        _layout_split(slide, w, h, structure, colors, image_bytes)
    elif layout == "two_column":
        _layout_two_column(slide, w, h, structure, colors)
    elif layout == "big_number":
        _layout_big_number(slide, w, h, structure, colors)
    else:
        _layout_bullets(slide, w, h, structure, colors)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


async def generate_slide_pptx(payload: dict[str, Any]) -> dict[str, Any]:
    brief = str(payload.get("brief") or "").strip()
    if not brief:
        raise ValueError("brief is required")
    title = str(payload.get("title") or "").strip()
    style_skill = str(payload.get("style_skill") or payload.get("styleSkill") or "")
    deck_context = payload.get("deck_context") or payload.get("deckContext")
    if deck_context is not None and not isinstance(deck_context, dict):
        deck_context = None
    images = payload.get("images") or []
    width = int(payload.get("width") or 1280)
    height = int(payload.get("height") or 720)

    structure = await _llm_structure(
        brief=brief, title=title, style_skill=style_skill, deck_context=deck_context
    )
    if title and not structure.get("title"):
        structure["title"] = title

    image_bytes = None
    role = str(structure.get("image_role") or "supporting").lower()
    if role != "none" and isinstance(images, list):
        for item in images[:3]:
            url = item.get("url") if isinstance(item, dict) else item
            if isinstance(url, str) and url.startswith("http"):
                image_bytes = await _fetch_image(url)
                if image_bytes:
                    break

    colors = _pick_colors(style_skill)
    pptx_bytes = _build_pptx(
        structure, colors=colors, image_bytes=image_bytes, width_px=width, height_px=height
    )
    _, _, model = _llm_openai_endpoint()
    model = model or os.environ.get("NVIDIA_DEFAULT_MODEL", "z-ai/glm-5.2")
    return {
        "pptx_base64": base64.b64encode(pptx_bytes).decode("ascii"),
        "model": f"trivena-local/{model}",
        "title": structure.get("title"),
        "layout": structure.get("layout"),
    }
